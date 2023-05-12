from pptx import Presentation
from pptx.util import Inches
import copy, screenshooter, os

def copy_slide_from_external_prs(prs):
    external_pres = Presentation(os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Assets/Template.pptx'))
    ext_slide = external_pres.slides[1]
    SLD_LAYOUT = 0
    slide_layout = prs.slide_layouts[SLD_LAYOUT]
    curr_slide = prs.slides.add_slide(slide_layout)
    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return prs

def move_slide(presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst 
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
        

def delete_slide(presentation,  index):
        xml_slides = presentation.slides._sldIdLst 
        slides = list(xml_slides)
        xml_slides.remove(slides[index])


def aqui(path):

    os.mkdir(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp")}')

    prs = Presentation(os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Assets/Template.pptx'))

    valid = screenshooter.take_screenshot()
    if valid:
        i = 0
        for _ in range(4):

            if i == 0:
                prs.slides[1+i].shapes.add_picture(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/front.png")}', Inches(1.9), Inches(1), width=Inches(7), height=Inches(5))
                os.remove(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/front.png")}')
            elif i == 1:
                prs.slides[1+i].shapes.add_picture(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/bottom.png")}', Inches(1.9), Inches(1), width=Inches(7), height=Inches(5))
                os.remove(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/bottom.png")}')
            elif i == 2:
                prs.slides[1+i].shapes.add_picture(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/right.png")}', Inches(0.1), Inches(1), width=Inches(5.25), height=Inches(5))
                os.remove(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/right.png")}')
                prs.slides[1+i].shapes.add_picture(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/left.png")}', Inches(5.5), Inches(1), width=Inches(5.25), height=Inches(5))
                os.remove(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp/left.png")}')

            prs = copy_slide_from_external_prs(prs)

            move_slide(prs, 2+i, 3+i)
            i+=1

        delete_slide(prs, 1+i)
        delete_slide(prs, i)

        os.rmdir(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp")}')

        try:
            os.remove(path)
        except FileNotFoundError:
            pass

        prs.save(path)

        return(True)
    
    else:
         
        os.rmdir(f'{os.path.join(os.path.dirname(os.path.abspath(__file__)), "Temp")}')

        return(False)