from pptx import Presentation
from pptx.util import Inches
import copy, screenshooter, os

def copy_slide_from_external_prs(prs):
    # copy from external presentation all objects into the existing presentation
    external_pres = Presentation("Fichiers/Pptx/ENR-0305 C Rapport 3D humain.pptx")
    # specify the slide you want to copy the contents from
    ext_slide = external_pres.slides[1]
    # Define the layout you want to use from your generated pptx
    SLD_LAYOUT = 0
    slide_layout = prs.slide_layouts[SLD_LAYOUT]
    # create now slide, to copy contents to 
    curr_slide = prs.slides.add_slide(slide_layout)
    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated
    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return prs

def move_slide(presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
        

def delete_slide(presentation,  index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[index])


prs = Presentation('Fichiers/Pptx/ENR-0305 C Rapport 3D humain.pptx')


# screenshooter.take_screenshot()

i = 0
for _ in range(5):

    if i == 0:
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/front.png', Inches(0.1), Inches(1), width=Inches(5), height=Inches(5))
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/front45.png', Inches(5.75), Inches(1), width=Inches(5), height=Inches(5))
    elif i == 1:
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/bottom.png', Inches(0.1), Inches(1), width=Inches(5), height=Inches(5))
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/bottom45.png', Inches(5.75), Inches(1), width=Inches(5), height=Inches(5))
    elif i == 2:
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/right.png', Inches(0.1), Inches(1), width=Inches(5), height=Inches(5))
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/left.png', Inches(5.75), Inches(1), width=Inches(5), height=Inches(5))
    elif i == 3:
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/right45.png', Inches(0.1), Inches(1), width=Inches(5), height=Inches(5))
        pic = prs.slides[1+i].shapes.add_picture(f'Fichiers/Images/left45.png', Inches(5.75), Inches(1), width=Inches(5), height=Inches(5))

    prs = copy_slide_from_external_prs(prs)

    move_slide(prs, 2+i, 3+i)
    i+=1

delete_slide(prs, 1+i)
delete_slide(prs, i)

try:
    os.remove("Fichiers/Pptx/test.pptx")
except FileNotFoundError:
    pass

prs.save('Fichiers/Pptx/test.pptx')